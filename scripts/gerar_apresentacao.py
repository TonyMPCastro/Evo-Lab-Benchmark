import os
from pptx import Presentation
from pptx.util import Inches, Pt, Cm
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# 1. Inicializa Apresentação
prs = Presentation()

# Proporção Widescreen 16:9
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# 2. Definição de Cores do Tema (Azul escuro, branco, cinza)
COLOR_DARK_BLUE = RGBColor(10, 30, 80)
COLOR_WHITE = RGBColor(255, 255, 255)
COLOR_GRAY = RGBColor(100, 100, 100)
COLOR_LIGHT_GRAY = RGBColor(240, 240, 240)

# Caminho para logotipo
LOGO_PATH = "logoufma.png"

def apply_background(slide):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = COLOR_WHITE

def add_footer(slide, slide_num):
    # Adicionar logotipo pequeno no rodapé
    if os.path.exists(LOGO_PATH):
        slide.shapes.add_picture(LOGO_PATH, Inches(0.2), Inches(6.8), height=Inches(0.5))
    
    # Texto de Rodapé
    txBox = slide.shapes.add_textbox(Inches(1.5), Inches(6.9), Inches(8), Inches(0.5))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "A. Castro | Análise de Algoritmos Bioinspirados em Problemas de Otimização"
    p.font.size = Pt(12)
    p.font.color.rgb = COLOR_GRAY
    
    # Número do slide
    txBox_num = slide.shapes.add_textbox(Inches(12.5), Inches(6.9), Inches(0.5), Inches(0.5))
    tf_num = txBox_num.text_frame
    p_num = tf_num.paragraphs[0]
    p_num.text = str(slide_num)
    p_num.font.size = Pt(12)
    p_num.font.color.rgb = COLOR_GRAY

def add_placeholder_box(slide, left, top, width, height, text):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    fill = shape.fill
    fill.solid()
    fill.fore_color.rgb = COLOR_LIGHT_GRAY
    line = shape.line
    line.color.rgb = COLOR_GRAY
    line.dash_style = 6 # Dashed
    
    tf = shape.text_frame
    tf.text = f"[ ESPAÇO RESERVADO PARA IMAGEM ]\n{text}"
    for p in tf.paragraphs:
        p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(14)
        p.font.color.rgb = COLOR_DARK_BLUE

def add_slide(title, bullets, notes, slide_num, placeholder_info=None):
    layout = prs.slide_layouts[1] # Title and Content
    slide = prs.slides.add_slide(layout)
    apply_background(slide)
    
    # Configurar Título
    title_shape = slide.shapes.title
    title_shape.text = title
    title_p = title_shape.text_frame.paragraphs[0]
    title_p.font.color.rgb = COLOR_DARK_BLUE
    title_p.font.size = Pt(40)
    title_p.font.bold = True
    
    # Configurar Corpo (Marcadores)
    body_shape = slide.placeholders[1]
    # Mudar tamanho do corpo dependendo se há placeholder ou não
    if placeholder_info:
        body_shape.width = Inches(5.5)
        
    tf = body_shape.text_frame
    tf.clear()
    
    for bullet in bullets:
        p = tf.add_paragraph()
        p.text = bullet
        p.font.size = Pt(24)
        p.font.color.rgb = COLOR_DARK_BLUE
        #p.level = 0
    
    # Adicionar Placeholder Box se fornecido
    if placeholder_info:
        # placeholder_info = {"left": ..., "top": ..., "width": ..., "height": ..., "text": ...}
        add_placeholder_box(
            slide, 
            placeholder_info['left'], 
            placeholder_info['top'], 
            placeholder_info['width'], 
            placeholder_info['height'], 
            placeholder_info['text']
        )
        
    # Adicionar Notas
    notes_slide = slide.notes_slide
    notes_tf = notes_slide.notes_text_frame
    notes_tf.text = notes
    
    # Adicionar Rodapé
    add_footer(slide, slide_num)
    
    return slide

# ================= SLIDE 1 =================
layout_title = prs.slide_layouts[0]
slide_1 = prs.slides.add_slide(layout_title)
apply_background(slide_1)

# Centralizar logo no topo, se houver
if os.path.exists(LOGO_PATH):
    slide_1.shapes.add_picture(LOGO_PATH, Inches(5.9), Inches(0.5), width=Inches(1.5))

title_shape = slide_1.shapes.title
title_shape.text = "ANÁLISE DE ALGORITMOS BIOINSPIRADOS EM PROBLEMAS DE OTIMIZAÇÃO"
title_shape.top = Inches(2.5)
title_p = title_shape.text_frame.paragraphs[0]
title_p.font.color.rgb = COLOR_DARK_BLUE
title_p.font.bold = True
title_p.font.size = Pt(36)

subtitle_shape = slide_1.placeholders[1]
subtitle_shape.top = Inches(4.5)
subtitle_tf = subtitle_shape.text_frame
subtitle_tf.text = "Autor: Antonio Marcos Patricio Castro\nUniversidade Federal do Maranhão (UFMA)\nCurso: Bacharelado Interdisciplinar em Ciência e Tecnologia\nOrientador: Prof. Dr. Pedro Baptista Fernandes\n2026"
for p in subtitle_tf.paragraphs:
    p.font.size = Pt(20)
    p.font.color.rgb = COLOR_GRAY

notes_slide = slide_1.notes_slide
notes_slide.notes_text_frame.text = "Falar: Apresentação do TCC e cumprimentos à banca examinadora.\nTempo estimado: 1 minuto.\nPontos importantes: Agradecer ao orientador."
# ============================================

slides_data = [
    {
        "title": "Agenda",
        "bullets": [
            "Contextualização e Motivação",
            "Problema de Pesquisa e Objetivos",
            "Fundamentação Teórica",
            "Algoritmos Estudados e Arquitetura",
            "Metodologia Experimental e Benchmarks",
            "Resultados e Análise Estatística",
            "Conclusão e Trabalhos Futuros"
        ],
        "notes": "Falar: Apresentar rapidamente os tópicos da defesa.\nTempo: 30 segundos.\nPontos importantes: Mostrar organização e fluidez.",
        "placeholder": None
    },
    {
        "title": "Contextualização",
        "bullets": [
            "Avanço em IA aumenta a complexidade matemática.",
            "Problemas NP-difíceis, não convexos e alta dimensão.",
            "Soluções clássicas aprisionadas em ótimos locais.",
            "Uso consolidado de meta-heurísticas bioinspiradas."
        ],
        "notes": "Falar: Explicar o crescimento de problemas de otimização em dados.\nTempo: 1 minuto.\nPontos: Modelos complexos exigem solvers robustos.\nPergunta possível: Por que métodos exatos falham? R: Porque a derivada não resolve espaços não-convexos.",
        "placeholder": {
            "left": Inches(7), "top": Inches(2.0), "width": Inches(5.5), "height": Inches(4.5),
            "text": "Inserir imagem de uma topologia não-convexa (ex: cheia de picos e vales) ilustrando Ótimos Locais."
        }
    },
    {
        "title": "Motivação",
        "bullets": [
            "Algoritmos originais carecem de poder de exploração.",
            "Híbridos trazem custo computacional extremamente elevado.",
            "Faltam avaliações padronizadas, reproduzíveis e estatísticas.",
            "Necessidade de validar arquiteturas frente à esparsidade."
        ],
        "notes": "Falar: Por que esse trabalho existe. Falta padronização empírica.\nTempo: 1 minuto.\nPontos: Mostrar que a literatura está confusa e enviesada.",
        "placeholder": None
    },
    {
        "title": "Problema de Pesquisa",
        "bullets": [
            "O incremento de complexidade justifica o alto custo?",
            "Como mitigar a dependência de parâmetros rígidos?",
            "Arquiteturas híbridas superam o dilema exploração vs explotação?",
            "Qual o impacto real da Maldição da Dimensionalidade?"
        ],
        "notes": "Falar: Explicar as perguntas centrais do trabalho.\nTempo: 1 minuto.\nPontos: A estagnação algorítmica e a rigidez paramétrica.",
        "placeholder": {
            "left": Inches(7), "top": Inches(2.0), "width": Inches(5.5), "height": Inches(4.5),
            "text": "Inserir Diagrama de Interrogação / Fluxo do Problema (Rigidez -> Estagnação -> Solução IA)."
        }
    },
    {
        "title": "Objetivos",
        "bullets": [
            "Geral: Analisar o comportamento de algoritmos bioinspirados unificadamente.",
            "Específico: Estruturar plataforma de benchmark padronizada.",
            "Específico: Avaliar complexidade e tempo de CPU.",
            "Específico: Aplicar em FS e HPO de Machine Learning.",
            "Específico: Validar com testes estatísticos não-paramétricos."
        ],
        "notes": "Falar: Apresentar o que entregaremos.\nTempo: 1 minuto.\nPontos: Ênfase na plataforma padronizada e no rigor estatístico.",
        "placeholder": None
    },
    {
        "title": "Fundamentação Teórica",
        "bullets": [
            "Teorema No Free Lunch: Nenhum otimizador é absoluto.",
            "Exploração: Varredura global e diversificação do enxame.",
            "Explotação: Refinamento intensivo da melhor solução encontrada.",
            "Evolução: Dos insetos sociais ao Aprendizado por Reforço."
        ],
        "notes": "Falar: Revisão teórica rápida.\nTempo: 1 minuto.\nPontos: O balanço entre explorar (descobrir) e explotar (refinar).",
        "placeholder": {
            "left": Inches(7), "top": Inches(2.0), "width": Inches(5.5), "height": Inches(4.5),
            "text": "Inserir diagrama de gangorra (Exploração vs Explotação)."
        }
    },
    {
        "title": "Algoritmos Estudados",
        "bullets": [
            "PSO: Pássaros. Rápido, mas sofre convergência prematura.",
            "ACO: Formigas e estigmergia. Excelente em grafos.",
            "GWO: Lobos e hierarquia. Alta precisão matemática.",
            "BWO: Belugas e Whale Fall. Colapsa em alta dimensão.",
            "QL-GWO: IA Tabular adaptando as fases dinamicamente.",
            "DRL-BWO: Redes Profundas para seleção binária acelerada."
        ],
        "notes": "Falar: Visão geral das seis técnicas testadas.\nTempo: 2 minutos.\nPontos: Contraste entre clássicos, modernos e híbridos (IA).",
        "placeholder": {
            "left": Inches(8.5), "top": Inches(2.0), "width": Inches(4.5), "height": Inches(4.5),
            "text": "Inserir Tabela Comparativa (Pros/Cons/Complexidade) dos Algoritmos."
        }
    },
    {
        "title": "Arquitetura da Plataforma",
        "bullets": [
            "Modelagem Orientada a Objetos para isolar viés.",
            "Classe de Controle: Orquestra chamadas estritamente justas.",
            "Classe Base: Padroniza contagem interna de NFEs.",
            "Módulo Gráfico: Renderização visual totalmente acadêmica."
        ],
        "notes": "Falar: Como a plataforma garante justiça e isenção.\nTempo: 1 minuto.\nPontos: O sistema impede vantagens algorítmicas de baixo nível.",
        "placeholder": {
            "left": Inches(7.5), "top": Inches(2.0), "width": Inches(5.0), "height": Inches(4.5),
            "text": "Redesenhar o Fluxograma do Trabalho original (Arquitetura POO)."
        }
    },
    {
        "title": "Fluxo Experimental",
        "bullets": [
            "Configuração inicial de agentes, dimensões e repetições.",
            "Submissão aos problemas matemáticos base.",
            "Extração de métricas operacionais (Custo e Tempo).",
            "Tratamento de outliers e execução de pós-testes estatísticos.",
            "Consolidação via Heatmaps, Violin Plots e Boxplots."
        ],
        "notes": "Falar: Passo a passo da metodologia.\nTempo: 1 minuto.\nPontos: As 30 execuções isolam fatores puramente aleatórios.",
        "placeholder": {
            "left": Inches(8), "top": Inches(2.0), "width": Inches(4.5), "height": Inches(4.5),
            "text": "Inserir diagrama de blocos: Config-> Problema-> Coleta-> Estatística-> Resultado."
        }
    },
    {
        "title": "Funções Benchmark",
        "bullets": [
            "Rastrigin: Armadilhas locais em formato de grade regular.",
            "Ackley, Sphere, Rosenbrock, Schwefel, Griewank.",
            "Avaliam escape sob extrema multimodalidade e esparsidade.",
            "Forçam dimensões rigorosas padronizadas em D=30 e D=50."
        ],
        "notes": "Falar: O campo de prova matemático.\nTempo: 1 minuto.\nPontos: Rastrigin é severa, testa capacidade exploratória bruta.",
        "placeholder": {
            "left": Inches(7), "top": Inches(2.0), "width": Inches(5.5), "height": Inches(4.5),
            "text": "Inserir pequenas imagens 3D das funções Rastrigin e Ackley em Tabela."
        }
    },
    {
        "title": "Problemas Reais de Machine Learning",
        "bullets": [
            "Feature Selection (FS): Altíssima esparsidade contínua (D=150).",
            "Hyperparameter Optimization (HPO): Espaço não estacionário (D=6).",
            "Base de Dados: GunPoint (séries temporais reais).",
            "Modelo Base: K-Nearest Neighbors (KNN)."
        ],
        "notes": "Falar: Aplicação prática na Engenharia de Dados.\nTempo: 1 minuto.\nPontos: Mostrar que não avaliamos apenas matemática teórica.",
        "placeholder": {
            "left": Inches(7), "top": Inches(2.0), "width": Inches(5.5), "height": Inches(4.5),
            "text": "Inserir representação esquemática do FS (Vetor de características 0/1)."
        }
    },
    {
        "title": "Metodologia Computacional",
        "bullets": [
            "Ambiente: Computação em nuvem via Google Colab.",
            "Linguagem Central: Python (Ecossistema científico).",
            "Modelagem Deep Learning: PyTorch e tensores isolados.",
            "Validação ML: Scikit-Learn com Validação Cruzada.",
            "Restrição Orçamentária Estrita (NFEs igualadas para todos)."
        ],
        "notes": "Falar: Stack tecnológico.\nTempo: 1 min.\nPergunta: Por que usar a nuvem? R: Isolar flutuações de background de CPU local.",
        "placeholder": None
    },
    {
        "title": "Resultados: Função Rastrigin",
        "bullets": [
            "PSO e ACO estagnaram nas armadilhas locais rapidamente.",
            "GWO mergulhou maciçamente à bacia global (Custo E-13).",
            "QL-GWO equilibrou bem a evasão de picos locais.",
            "Observado claro Trade-off entre Custo Computacional e Precisão."
        ],
        "notes": "Falar: Os resultados da matemática multivariada.\nTempo: 1.5 minuto.\nPontos: Mostrar o Gráfico de Convergência (Eixo Y logarítmico).",
        "placeholder": {
            "left": Inches(6.5), "top": Inches(1.5), "width": Inches(6.0), "height": Inches(5.0),
            "text": "Inserir Gráfico de Convergência da Função Rastrigin (Tempo x Custo)."
        }
    },
    {
        "title": "Resultados: Feature Selection",
        "bullets": [
            "Maldição da Dimensionalidade (D=150) puniu modelos simples.",
            "BWO colapsou; Whale Fall foi desastroso neste domínio.",
            "GWO exibiu robustez geométrica excepcional.",
            "QL-GWO seguiu de perto, com forte inteligência adaptativa."
        ],
        "notes": "Falar: O desastre de alguns algoritmos na vida real.\nTempo: 1 minuto.\nPontos: O algoritmo da Beluga falhou pela aleatoriedade cega.",
        "placeholder": {
            "left": Inches(7), "top": Inches(2.0), "width": Inches(5.5), "height": Inches(4.5),
            "text": "Inserir gráfico de barras ou boxplot do Erro de Validação FS."
        }
    },
    {
        "title": "Resultados Estatísticos",
        "bullets": [
            "Kruskal-Wallis e Friedman confirmam assimetria global p<0.05.",
            "Ranking Friedman consolida GWO e QL-GWO no topo.",
            "Mann-Whitney U comprova vitória matemática estrita.",
            "Cliff's Delta aferiu magnitude Absoluta (Grande) para GWO."
        ],
        "notes": "Falar: A prova cabal e científica.\nTempo: 1 minuto.\nPontos: Não dependemos da sorte, a estatística cravou a vitória.",
        "placeholder": {
            "left": Inches(7), "top": Inches(2.0), "width": Inches(5.5), "height": Inches(4.5),
            "text": "Inserir Mapa de Calor (Heatmap) do Cliff's Delta ou Ranking de Friedman."
        }
    },
    {
        "title": "Discussão e Desempenho Global",
        "bullets": [
            "QL-GWO provou ser o mais consistente (menor desvio).",
            "Q-Table substitui decaimentos cegos por micro-ajustes exatos.",
            "BWO possui arquitetura leve, mas inviável por baixa precisão.",
            "GWO lidera entre heurísticas puramente biológicas e matemáticas."
        ],
        "notes": "Falar: Quem venceu afinal?\nTempo: 1.5 minuto.\nPontos: QL-GWO venceu na vida real, GWO venceu na matemática teórica.",
        "placeholder": {
            "left": Inches(7.5), "top": Inches(2.0), "width": Inches(5.0), "height": Inches(4.5),
            "text": "Inserir Gráfico de Dispersão (Trade-Off Tempo CPU vs Custo Final)."
        }
    },
    {
        "title": "Contribuições da Pesquisa",
        "bullets": [
            "Plataforma algorítmica modular pronta para novas avaliações.",
            "Comprovação empírica da eficácia das hibridizações IA-Enxame.",
            "Refutação técnica do BWO para alta complexidade esparsa.",
            "Validação rigorosa que suprime vieses de hardware e sorte."
        ],
        "notes": "Falar: O que entregamos à sociedade e à universidade.\nTempo: 1 minuto.\nPontos: Destacar a plataforma de software criada.",
        "placeholder": None
    },
    {
        "title": "Conclusão",
        "bullets": [
            "Objetivos metodológicos cumpridos com extrema precisão.",
            "O decaimento paramétrico rígido demonstrou-se obsoleto.",
            "QL-GWO desponta como escolha mandatória para engenharia.",
            "GWO sustenta soberania nas heurísticas estritamente matemáticas."
        ],
        "notes": "Falar: Fechamento das análises.\nTempo: 1 minuto.\nPontos: Responder à pergunta de pesquisa (Sim, a IA híbrida vale a pena).",
        "placeholder": None
    },
    {
        "title": "Limitações",
        "bullets": [
            "Custo computacional elevado restringe avaliação massiva (NFE).",
            "Impossibilidade de treinar Redes Neurais Profundas (LSTMs).",
            "Matriz tabular (Q-Learning) exposta à explosão combinatória.",
            "Restrito a 30 execuções dada a infraestrutura disponível."
        ],
        "notes": "Falar: O que não foi possível fazer.\nTempo: 1 min.\nPergunta: Por que 30 execuções? R: Padrão mínimo para Teorema do Limite Central.",
        "placeholder": None
    },
    {
        "title": "Trabalhos Futuros",
        "bullets": [
            "Sintonia hiperparamétrica para Deep Learning (CNN, LSTM).",
            "Otimização de rotas em ambientes não-estacionários e drones.",
            "Substituição do QL Tabular por DDPG (Atores Contínuos).",
            "Implementação de processamento paralelo e paralelismo de GPU."
        ],
        "notes": "Falar: Para onde vai o projeto a partir de hoje.\nTempo: 1 minuto.\nPontos: As melhorias óbvias da plataforma.",
        "placeholder": None
    },
    {
        "title": "Agradecimentos",
        "bullets": [
            "Agradeço ao orientador, familiares e à UFMA.",
            "Obrigado pela atenção de todos.",
            "Aberto a perguntas e arguições da banca."
        ],
        "notes": "Falar: Agradecimentos finais à banca.\nTempo: 30 segundos.",
        "placeholder": None
    }
]

# Adicionar os 20 slides restantes
for i, data in enumerate(slides_data):
    slide_num = i + 2
    add_slide(data["title"], data["bullets"], data["notes"], slide_num, data.get("placeholder"))

# Salvar Arquivo
prs.save("Apresentacao_TCC_Antonio_Castro.pptx")
print("Apresentação gerada com sucesso: Apresentacao_TCC_Antonio_Castro.pptx")
